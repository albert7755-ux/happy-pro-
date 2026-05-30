"""
8-Page Private Banking Proposal Generator
Style: Morgan Stanley / BlackRock / UBS Deep Navy & Gold
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import pptx.oxml.ns as nsmap
from lxml import etree
import math

# ── Color Palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy
NAVY2     = RGBColor(0x12, 0x28, 0x58)
GOLD      = RGBColor(0xC8, 0xA8, 0x4B)
GOLD2     = RGBColor(0xE8, 0xCB, 0x7A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTGRAY = RGBColor(0xF4, 0xF6, 0xFA)
MIDGRAY   = RGBColor(0xB0, 0xB8, 0xCC)
DARKGRAY  = RGBColor(0x44, 0x4C, 0x60)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helper utilities ───────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb

def add_divider(slide, l, t, w, color=GOLD, thickness=Pt(2)):
    line = slide.shapes.add_shape(1, l, t, w, int(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def navy_header(slide, title, subtitle=""):
    """Full-width navy top bar with gold accent line."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_divider(slide, 0, Inches(1.15), SLIDE_W, color=GOLD, thickness=Pt(3))
    add_text(slide, title, Inches(0.4), Inches(0.15),
             Inches(10), Inches(0.55),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72),
                 Inches(10), Inches(0.35),
                 font_size=11, color=MIDGRAY, align=PP_ALIGN.LEFT)

def footer(slide, page_num, total=8, firm="私人財富管理部門"):
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), fill_color=NAVY)
    add_text(slide, f"© 2025 {firm}  |  本文件僅供內部提案使用，請勿外流",
             Inches(0.3), SLIDE_H - Inches(0.34),
             Inches(10), Inches(0.3),
             font_size=7.5, color=MIDGRAY, align=PP_ALIGN.LEFT)
    add_text(slide, f"{page_num} / {total}",
             Inches(12.5), SLIDE_H - Inches(0.34),
             Inches(0.7), Inches(0.3),
             font_size=8, color=GOLD, align=PP_ALIGN.RIGHT)

def kpi_box(slide, l, t, w, h, value, label, value_color=GOLD):
    add_rect(slide, l, t, w, h, fill_color=NAVY2)
    add_rect(slide, l, t, Inches(0.05), h, fill_color=GOLD)
    add_text(slide, value, l + Inches(0.12), t + Inches(0.1),
             w - Inches(0.15), Inches(0.42),
             font_size=24, bold=True, color=value_color, align=PP_ALIGN.LEFT)
    add_text(slide, label, l + Inches(0.12), t + Inches(0.52),
             w - Inches(0.15), Inches(0.28),
             font_size=9, color=MIDGRAY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 – COVER
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# Full navy background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
# Gold accent bar left edge
add_rect(s1, 0, 0, Inches(0.25), SLIDE_H, fill_color=GOLD)
# Gold horizontal accent
add_rect(s1, Inches(0.25), Inches(3.6), Inches(8), Pt(3), fill_color=GOLD)

# Logo placeholder (top-right)
logo_box = add_rect(s1, SLIDE_W - Inches(2.8), Inches(0.25),
                    Inches(2.5), Inches(0.7), fill_color=NAVY2)
add_text(s1, "PRIVATE WEALTH", SLIDE_W - Inches(2.8), Inches(0.35),
         Inches(2.5), Inches(0.5),
         font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Main title
add_text(s1, "投資方案提案報告",
         Inches(0.6), Inches(1.8), Inches(9), Inches(0.9),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "INVESTMENT PROPOSAL",
         Inches(0.6), Inches(2.7), Inches(9), Inches(0.5),
         font_size=18, bold=False, color=GOLD2, align=PP_ALIGN.LEFT, italic=True)

# Subtitle block
add_text(s1, "客戶姓名：陳○○ 先生 / 女士",
         Inches(0.6), Inches(3.85), Inches(6), Inches(0.35),
         font_size=13, color=MIDGRAY)
add_text(s1, "提案日期：2025年6月",
         Inches(0.6), Inches(4.2), Inches(6), Inches(0.35),
         font_size=13, color=MIDGRAY)
add_text(s1, "理財顧問：私人財富管理部門",
         Inches(0.6), Inches(4.55), Inches(6), Inches(0.35),
         font_size=13, color=MIDGRAY)

# Bottom tagline
add_text(s1, "專業．信賴．長期財富傳承",
         Inches(0.6), Inches(6.5), Inches(8), Inches(0.4),
         font_size=11, color=GOLD, italic=True)

footer(s1, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 – 客戶現況
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHTGRAY)
navy_header(s2, "客戶現況分析", "Client Profile & Current Situation")
footer(s2, 2)

# Profile card left
add_rect(s2, Inches(0.3), Inches(1.4), Inches(3.8), Inches(5.5), fill_color=WHITE)
add_rect(s2, Inches(0.3), Inches(1.4), Inches(3.8), Pt(4), fill_color=GOLD)
add_text(s2, "基本資料", Inches(0.5), Inches(1.5), Inches(3.4), Inches(0.4),
         font_size=13, bold=True, color=NAVY)

profile_items = [
    ("投資目標", "穩定現金流 + 資產增值"),
    ("風險屬性", "穩健型（中低風險）"),
    ("投資期間", "5 – 10 年"),
    ("可投資資產", "NT$ 5,000 萬"),
    ("稅務身份", "中華民國居民"),
    ("流動性需求", "每月配息 NT$15萬+"),
    ("特殊考量", "退休規劃、財富傳承"),
]
for i, (k, v) in enumerate(profile_items):
    y = Inches(2.1) + i * Inches(0.6)
    add_rect(s2, Inches(0.35), y, Inches(3.7), Inches(0.55),
             fill_color=LIGHTGRAY if i % 2 == 0 else WHITE)
    add_text(s2, k, Inches(0.45), y + Inches(0.1), Inches(1.3), Inches(0.35),
             font_size=9.5, bold=True, color=DARKGRAY)
    add_text(s2, v, Inches(1.8), y + Inches(0.1), Inches(2.1), Inches(0.35),
             font_size=9.5, color=NAVY)

# Current portfolio status
add_rect(s2, Inches(4.4), Inches(1.4), Inches(8.6), Inches(2.5), fill_color=WHITE)
add_rect(s2, Inches(4.4), Inches(1.4), Inches(8.6), Pt(4), fill_color=GOLD)
add_text(s2, "現有資產配置概況", Inches(4.6), Inches(1.5), Inches(6), Inches(0.4),
         font_size=13, bold=True, color=NAVY)

current_assets = [
    ("定存 / 貨幣市場", "NT$1,500萬", "30%", "1.8%"),
    ("台股基金", "NT$800萬", "16%", "不穩定"),
    ("外幣定存 (USD)", "NT$1,200萬", "24%", "4.5%"),
    ("保險儲蓄", "NT$1,000萬", "20%", "2.2%"),
    ("其他 (房產/黃金)", "NT$500萬", "10%", "—"),
]
headers = ["類別", "金額", "佔比", "預估年收益率"]
for ci, h in enumerate(headers):
    x = Inches(4.5) + ci * Inches(2.0)
    add_rect(s2, x, Inches(2.0), Inches(2.0), Inches(0.35), fill_color=NAVY)
    add_text(s2, h, x, Inches(2.0), Inches(2.0), Inches(0.35),
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(current_assets):
    bg = LIGHTGRAY if ri % 2 == 0 else WHITE
    for ci, val in enumerate(row):
        x = Inches(4.5) + ci * Inches(2.0)
        y = Inches(2.35) + ri * Inches(0.32)
        add_rect(s2, x, y, Inches(2.0), Inches(0.32), fill_color=bg)
        add_text(s2, val, x + Inches(0.05), y + Inches(0.04),
                 Inches(1.9), Inches(0.25),
                 font_size=8.5, color=NAVY, align=PP_ALIGN.CENTER)

# Pain points
add_rect(s2, Inches(4.4), Inches(4.1), Inches(8.6), Inches(2.8), fill_color=WHITE)
add_rect(s2, Inches(4.4), Inches(4.1), Inches(8.6), Pt(4), fill_color=GOLD)
add_text(s2, "核心需求與痛點", Inches(4.6), Inches(4.2), Inches(6), Inches(0.4),
         font_size=13, bold=True, color=NAVY)
pain_points = [
    ("①", "現金流不足", "現有配息收入低於退休月支出目標 NT$15萬"),
    ("②", "集中風險高", "資產過度集中台幣定存與台股，匯率與市場風險大"),
    ("③", "報酬率偏低", "整體資產年化報酬率約 2.5%，未跟上通膨"),
    ("④", "全球分散不足", "缺乏美元優質公司債與全球多元資產配置"),
]
for i, (num, title, desc) in enumerate(pain_points):
    y = Inches(4.7) + i * Inches(0.52)
    add_text(s2, num, Inches(4.55), y, Inches(0.3), Inches(0.4),
             font_size=14, bold=True, color=GOLD)
    add_text(s2, title, Inches(4.85), y, Inches(2.0), Inches(0.4),
             font_size=10, bold=True, color=NAVY)
    add_text(s2, desc, Inches(6.9), y, Inches(5.8), Inches(0.4),
             font_size=9, color=DARKGRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 – 市場觀點
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHTGRAY)
navy_header(s3, "2025 市場觀點與投資主題", "Market Outlook & Investment Themes")
footer(s3, 3)

themes = [
    ("🔵", "投資等級美元公司債",
     "聯準會降息循環啟動，長債殖利率仍維持 4.8-5.6% 吸引力。\n優質企業信用穩健，提供穩定現金流與資本利得機會。"),
    ("🟡", "多元資產收益基金",
     "全球股債混合策略，年化配息率約 7-10%。\n透過動態資產配置降低波動、提升風險調整後報酬。"),
    ("🟢", "新興市場精選債券",
     "亞洲優質主權債與金融債，殖利率溢價吸引人。\n選擇性布局印度、沙烏地主權債，分散地緣風險。"),
    ("🔴", "台幣/外幣資產再平衡",
     "台幣升值壓力下，適度增加外幣資產至 50%+。\n美元資產提供天然避險並享受較高利率環境。"),
]
for i, (icon, title, desc) in enumerate(themes):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * Inches(6.55)
    y = Inches(1.5) + row * Inches(2.7)
    add_rect(s3, x, y, Inches(6.2), Inches(2.55), fill_color=WHITE)
    add_rect(s3, x, y, Inches(0.08), Inches(2.55), fill_color=GOLD)
    add_text(s3, icon + "  " + title, x + Inches(0.2), y + Inches(0.12),
             Inches(5.8), Inches(0.5),
             font_size=13, bold=True, color=NAVY)
    add_divider(s3, x + Inches(0.2), y + Inches(0.68), Inches(5.6), color=GOLD, thickness=Pt(1))
    add_text(s3, desc, x + Inches(0.2), y + Inches(0.82),
             Inches(5.8), Inches(1.55),
             font_size=10, color=DARKGRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 – 資產配置圓餅圖（用形狀模擬）
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHTGRAY)
navy_header(s4, "建議資產配置", "Recommended Asset Allocation")
footer(s4, 4)

alloc = [
    ("投資等級公司債", 40, NAVY,  GOLD),
    ("多元收益基金",   25, RGBColor(0x1A,0x4A,0x8C), GOLD2),
    ("新興市場債券",   15, RGBColor(0x2E,0x6B,0xAD), WHITE),
    ("全球股票基金",   12, RGBColor(0x44,0x88,0xC8), WHITE),
    ("現金/貨幣市場",  8,  MIDGRAY, NAVY),
]

# Draw pie using arc shapes (approximated with rectangles + text)
cx, cy, r = Inches(4.5), Inches(4.2), Inches(2.3)

# Legend table (right side)
add_rect(s4, Inches(8.5), Inches(1.5), Inches(4.5), Inches(5.5), fill_color=WHITE)
add_rect(s4, Inches(8.5), Inches(1.5), Inches(4.5), Pt(3), fill_color=GOLD)
add_text(s4, "配置明細", Inches(8.7), Inches(1.6), Inches(4), Inches(0.4),
         font_size=13, bold=True, color=NAVY)

legend_colors = [NAVY, RGBColor(0x1A,0x4A,0x8C), RGBColor(0x2E,0x6B,0xAD),
                 RGBColor(0x44,0x88,0xC8), MIDGRAY]
for i, (name, pct, clr, txt_clr) in enumerate(alloc):
    y = Inches(2.2) + i * Inches(0.9)
    add_rect(s4, Inches(8.6), y, Inches(0.25), Inches(0.55), fill_color=clr)
    add_text(s4, name, Inches(8.95), y + Inches(0.05),
             Inches(2.8), Inches(0.3),
             font_size=11, bold=True, color=NAVY)
    add_text(s4, f"{pct}%", Inches(8.95), y + Inches(0.32),
             Inches(2.8), Inches(0.25),
             font_size=10, color=DARKGRAY)
    add_text(s4, f"NT${pct*50:,}萬", Inches(11.2), y + Inches(0.15),
             Inches(1.5), Inches(0.3),
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

# Pie chart using python-pptx chart
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

cd = ChartData()
cd.categories = [a[0] for a in alloc]
cd.add_series("配置", [a[1] for a in alloc])

chart_placeholder = s4.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT,
    Inches(0.3), Inches(1.4),
    Inches(7.8), Inches(5.7),
    cd
)
chart = chart_placeholder.chart
chart.has_legend = False
chart.has_title  = False

# Color the slices
from pptx.oxml.ns import qn
plot = chart.plots[0]
ser  = plot.series[0]
pts  = ser._element.findall(qn('c:dPt'))

pie_colors = [NAVY, RGBColor(0x1A,0x4A,0x8C), RGBColor(0x2E,0x6B,0xAD),
              RGBColor(0x44,0x88,0xC8), MIDGRAY]

# Add center label
add_text(s4, "NT$\n5,000萬", cx - Inches(0.8), cy - Inches(0.5),
         Inches(1.6), Inches(1.0),
         font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# KPI row below chart
kpis = [("預期年化報酬", "6.8%"), ("預期年化波動", "5.2%"), ("夏普比率", "0.94"), ("年化配息收入", "NT$288萬")]
for i, (lbl, val) in enumerate(kpis):
    x = Inches(0.3) + i * Inches(3.25)
    kpi_box(s4, x, Inches(7.0) - Inches(0.6), Inches(3.0), Inches(0.82), val, lbl)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 – 商品比較
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHTGRAY)
navy_header(s5, "核心商品比較", "Core Product Comparison")
footer(s5, 5)

products = [
    ("商品名稱",      "Alphabet公司債\n5.5% 2046", "Meta公司債\n5.5% 2045", "PIMCO收益增長基金", "聯博美國收益基金", "施羅德環球收息"),
    ("類型",          "投資等級公司債", "投資等級公司債", "多元收益基金", "美元收益基金", "全球債券基金"),
    ("殖利率/配息率", "5.21%",          "5.46%",          "9.00%",        "7.81%",        "5.51%"),
    ("到期/存續期",   "2046年",         "2045年",         "—",            "—",            "—"),
    ("幣別",          "美元",           "美元",           "美元",         "美元",         "美元"),
    ("信用評級",      "AA",             "AA",             "投資等級",     "多元",         "投資等級"),
    ("配息頻率",      "半年配",         "半年配",         "月配",         "月配",         "月配"),
    ("建議配置",      "15%",            "15%",            "15%",          "10%",          "10%"),
    ("建議金額",      "NT$750萬",       "NT$750萬",       "NT$750萬",     "NT$500萬",     "NT$500萬"),
]

col_widths = [Inches(1.8), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)]
col_x = [Inches(0.15)]
for cw in col_widths[:-1]:
    col_x.append(col_x[-1] + cw)

for ri, row in enumerate(products):
    for ci, cell in enumerate(row):
        x = col_x[ci]
        y = Inches(1.3) + ri * Inches(0.57)
        w = col_widths[ci]
        if ri == 0:
            bg = NAVY
            fc = WHITE
            fsize = 9
            fb = True
        elif ci == 0:
            bg = NAVY2
            fc = GOLD
            fsize = 9
            fb = True
        else:
            bg = LIGHTGRAY if ri % 2 == 0 else WHITE
            fc = NAVY
            fsize = 9
            fb = False
        add_rect(s5, x, y, w, Inches(0.55), fill_color=bg)
        add_text(s5, cell, x + Inches(0.05), y + Inches(0.05),
                 w - Inches(0.1), Inches(0.47),
                 font_size=fsize, bold=fb, color=fc, align=PP_ALIGN.CENTER)

# Highlight best pick
add_rect(s5, col_x[1], Inches(1.3), col_widths[1] * 2, Inches(0.04), fill_color=GOLD)

add_text(s5, "★ 重點推薦商品",
         Inches(0.15), Inches(6.6), Inches(5), Inches(0.35),
         font_size=9.5, bold=True, color=GOLD)
add_text(s5, "以上殖利率與配息率資料來源：公司內部資料庫（2025/05），僅供參考，實際以各機構公告為準。",
         Inches(0.15), Inches(6.92), Inches(12), Inches(0.3),
         font_size=7.5, color=MIDGRAY, italic=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 – 現金流圖表
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHTGRAY)
navy_header(s6, "配息現金流試算", "Cash Flow Projection")
footer(s6, 6)

# KPI strip
kpi_data = [("年化配息收入", "NT$288萬"), ("平均殖利率", "5.76%"), ("月均入帳", "NT$24萬"), ("超過目標", "+NT$9萬/月")]
for i, (lbl, val) in enumerate(kpi_data):
    x = Inches(0.3) + i * Inches(3.2)
    kpi_box(s6, x, Inches(1.35), Inches(3.0), Inches(0.88), val, lbl)

# Monthly bar chart
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

months = ["1月","2月","3月","4月","5月","6月","7月","8月","9月","10月","11月","12月"]
# Simulated monthly cash flow (bonds pay biannually, funds monthly)
monthly_vals = [18.5, 27.3, 24.0, 28.8, 18.5, 34.2, 24.0, 22.1, 35.6, 24.0, 18.5, 33.8]

cd2 = ChartData()
cd2.categories = months
cd2.add_series("每月配息 (萬NTD)", monthly_vals)

chart2 = s6.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.3), Inches(2.4),
    Inches(8.0), Inches(4.5),
    cd2
).chart
chart2.has_legend = False
chart2.has_title  = True
chart2.chart_title.text_frame.text = "逐月配息入帳預測 (萬 NTD)"

# Detail table (right side)
add_rect(s6, Inches(8.7), Inches(2.4), Inches(4.3), Inches(4.5), fill_color=WHITE)
add_rect(s6, Inches(8.7), Inches(2.4), Inches(4.3), Pt(3), fill_color=GOLD)
add_text(s6, "各標的年化配息", Inches(8.9), Inches(2.5), Inches(3.8), Inches(0.38),
         font_size=11, bold=True, color=NAVY)

detail_rows = [
    ("Alphabet債 5.5%", "15%", "NT$67.5萬"),
    ("Meta債 5.5%",     "15%", "NT$67.5萬"),
    ("PIMCO收益增長",   "15%", "NT$67.5萬"),
    ("聯博美國收益",    "10%", "NT$39.1萬"),
    ("施羅德環球收息",  "10%", "NT$27.6萬"),
    ("其他配置",        "35%", "NT$18.8萬"),
]
hdrs = ["標的", "配置", "年配息"]
for ci, h in enumerate(hdrs):
    x = Inches(8.75) + ci * Inches(1.35)
    add_rect(s6, x, Inches(2.95), Inches(1.35), Inches(0.32), fill_color=NAVY)
    add_text(s6, h, x, Inches(2.95), Inches(1.35), Inches(0.32),
             font_size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(detail_rows):
    bg = LIGHTGRAY if ri % 2 == 0 else WHITE
    for ci, val in enumerate(row):
        x = Inches(8.75) + ci * Inches(1.35)
        y = Inches(3.27) + ri * Inches(0.52)
        add_rect(s6, x, y, Inches(1.35), Inches(0.5), fill_color=bg)
        fc = GOLD if ci == 2 else NAVY
        add_text(s6, val, x, y + Inches(0.08), Inches(1.35), Inches(0.36),
                 font_size=8.5, color=fc, align=PP_ALIGN.CENTER)

add_text(s6, "合計年配息",
         Inches(8.75), Inches(6.4), Inches(2.7), Inches(0.35),
         font_size=10, bold=True, color=NAVY)
add_text(s6, "NT$288萬",
         Inches(11.2), Inches(6.4), Inches(1.6), Inches(0.35),
         font_size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 – 預期成果
# ═══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHTGRAY)
navy_header(s7, "預期成果與績效分析", "Expected Outcomes & Performance")
footer(s7, 7)

# Before vs After comparison
add_rect(s7, Inches(0.3), Inches(1.35), Inches(6.0), Inches(2.9), fill_color=WHITE)
add_rect(s7, Inches(0.3), Inches(1.35), Inches(6.0), Pt(4), fill_color=MIDGRAY)
add_text(s7, "📌 目前配置", Inches(0.5), Inches(1.42), Inches(5.5), Inches(0.45),
         font_size=13, bold=True, color=DARKGRAY)

before = [("年化報酬率", "≈ 2.5%"), ("年化波動率", "≈ 4.0%"),
          ("月均配息", "NT$0 – 5萬"), ("美元資產比重", "24%"),
          ("全球分散程度", "低")]
for i, (k, v) in enumerate(before):
    y = Inches(1.95) + i * Inches(0.42)
    add_rect(s7, Inches(0.35), y, Inches(5.9), Inches(0.4),
             fill_color=LIGHTGRAY if i % 2 == 0 else WHITE)
    add_text(s7, k, Inches(0.45), y + Inches(0.06), Inches(3), Inches(0.3),
             font_size=9.5, color=DARKGRAY)
    add_text(s7, v, Inches(3.5), y + Inches(0.06), Inches(2.5), Inches(0.3),
             font_size=9.5, bold=True, color=DARKGRAY, align=PP_ALIGN.RIGHT)

add_rect(s7, Inches(6.7), Inches(1.35), Inches(6.3), Inches(2.9), fill_color=WHITE)
add_rect(s7, Inches(6.7), Inches(1.35), Inches(6.3), Pt(4), fill_color=GOLD)
add_text(s7, "✅ 建議配置後", Inches(6.9), Inches(1.42), Inches(5.8), Inches(0.45),
         font_size=13, bold=True, color=NAVY)

after = [("年化報酬率", "≈ 6.8%"), ("年化波動率", "≈ 5.2%"),
         ("月均配息", "NT$24萬+"), ("美元資產比重", "55%"),
         ("全球分散程度", "高")]
improvement = ["▲ +4.3pp", "▲ +1.2pp", "▲ +NT$19萬", "▲ +31pp", "大幅提升"]
for i, ((k, v), imp) in enumerate(zip(after, improvement)):
    y = Inches(1.95) + i * Inches(0.42)
    add_rect(s7, Inches(6.75), y, Inches(6.2), Inches(0.4),
             fill_color=LIGHTGRAY if i % 2 == 0 else WHITE)
    add_text(s7, k, Inches(6.85), y + Inches(0.06), Inches(2.5), Inches(0.3),
             font_size=9.5, color=NAVY)
    add_text(s7, v, Inches(9.4), y + Inches(0.06), Inches(1.5), Inches(0.3),
             font_size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s7, imp, Inches(10.9), y + Inches(0.06), Inches(1.8), Inches(0.3),
             font_size=9, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

# Arrow between
add_text(s7, "→", Inches(6.1), Inches(2.55), Inches(0.6), Inches(0.6),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# 5-year projection table
add_rect(s7, Inches(0.3), Inches(4.45), Inches(12.7), Inches(2.65), fill_color=WHITE)
add_rect(s7, Inches(0.3), Inches(4.45), Inches(12.7), Pt(3), fill_color=GOLD)
add_text(s7, "5年資產成長預測（含配息再投入）", Inches(0.5), Inches(4.52),
         Inches(8), Inches(0.4), font_size=12, bold=True, color=NAVY)

proj_headers = ["", "第1年末", "第2年末", "第3年末", "第4年末", "第5年末"]
proj_data = [
    ("保守情境 (4.5%)",  "5,225萬", "5,461萬", "5,707萬", "5,964萬", "6,232萬"),
    ("基準情境 (6.8%)",  "5,340萬", "5,703萬", "6,091萬", "6,505萬", "6,948萬"),
    ("樂觀情境 (9.0%)",  "5,450萬", "5,941萬", "6,475萬", "7,057萬", "7,692萬"),
]

col_w2 = Inches(12.5) / 6
for ci, h in enumerate(proj_headers):
    x = Inches(0.35) + ci * col_w2
    add_rect(s7, x, Inches(5.0), col_w2, Inches(0.38), fill_color=NAVY)
    add_text(s7, h, x, Inches(5.0), col_w2, Inches(0.38),
             font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(proj_data):
    for ci, val in enumerate(row):
        x = Inches(0.35) + ci * col_w2
        y = Inches(5.38) + ri * Inches(0.5)
        bg = LIGHTGRAY if ri % 2 == 0 else WHITE
        if ci == 0:
            bg = NAVY2
        add_rect(s7, x, y, col_w2, Inches(0.48), fill_color=bg)
        fc = GOLD if ci == 0 else (GOLD if ri == 1 else NAVY)
        add_text(s7, val, x, y + Inches(0.08), col_w2, Inches(0.34),
                 font_size=9, bold=(ci == 0), color=fc, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 – CLOSING / NEXT STEPS
# ═══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
add_rect(s8, 0, 0, Inches(0.25), SLIDE_H, fill_color=GOLD)

# Gold line
add_rect(s8, Inches(0.25), Inches(2.8), Inches(9), Pt(2), fill_color=GOLD)

add_text(s8, "下一步行動計劃", Inches(0.55), Inches(1.0),
         Inches(10), Inches(0.7),
         font_size=34, bold=True, color=WHITE)
add_text(s8, "Next Steps & Implementation Timeline",
         Inches(0.55), Inches(1.75), Inches(10), Inches(0.4),
         font_size=14, color=GOLD2, italic=True)

steps = [
    ("STEP 1", "確認投資意向書", "簽署 KYC 更新文件及投資意向確認書（預計 1 週內完成）"),
    ("STEP 2", "完成開戶作業",   "外匯帳戶開立 + 美元資金匯兌準備（預計 2 週）"),
    ("STEP 3", "分批建倉執行",   "依市場時機分 2-3 批購入各標的，降低時間點風險"),
    ("STEP 4", "定期績效回顧",   "每季提供配息報告 + 每半年進行資產配置再平衡"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.55) + (i % 2) * Inches(6.3)
    y = Inches(3.1) + (i // 2) * Inches(1.85)
    add_rect(s8, x, y, Inches(5.8), Inches(1.65), fill_color=NAVY2)
    add_rect(s8, x, y, Inches(5.8), Pt(3), fill_color=GOLD)
    add_text(s8, num, x + Inches(0.15), y + Inches(0.1),
             Inches(1.2), Inches(0.45),
             font_size=11, bold=True, color=GOLD)
    add_text(s8, title, x + Inches(0.15), y + Inches(0.55),
             Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=WHITE)
    add_text(s8, desc, x + Inches(0.15), y + Inches(1.0),
             Inches(5.5), Inches(0.55),
             font_size=9.5, color=MIDGRAY)

# Contact block
add_rect(s8, Inches(9.2), Inches(1.0), Inches(3.9), Inches(5.8), fill_color=NAVY2)
add_rect(s8, Inches(9.2), Inches(1.0), Inches(3.9), Pt(3), fill_color=GOLD)
add_text(s8, "聯絡資訊", Inches(9.4), Inches(1.1), Inches(3.5), Inches(0.4),
         font_size=12, bold=True, color=GOLD)
contact_lines = [
    "理財顧問",
    "王○○ 副總裁",
    "",
    "📞  02-2xxx-xxxx",
    "📧  advisor@bank.com.tw",
    "",
    "私人財富管理部門",
    "台北市信義區○○路100號",
    "",
    "服務時間",
    "週一至週五",
    "09:00 – 18:00",
]
for i, line in enumerate(contact_lines):
    bold = line in ["理財顧問", "王○○ 副總裁", "私人財富管理部門", "服務時間"]
    color = GOLD if bold else MIDGRAY
    add_text(s8, line, Inches(9.4), Inches(1.65) + i * Inches(0.33),
             Inches(3.5), Inches(0.32),
             font_size=10 if bold else 9, bold=bold, color=color)

footer(s8, 8)

# ── Save ───────────────────────────────────────────────────────
out_path = "/home/user/happy-pro-/私銀投資提案_2025.pptx"
prs.save(out_path)
print(f"✅ PPTX saved: {out_path}")
